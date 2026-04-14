"""
Build the 5-minute presentation as a .pptx file.

Can run in two environments:

  1. Locally on the repo (here) — finds pipeline_architecture.png and
     package_structure.png in the repo, uses placeholder boxes for
     figures that live on Google Drive. Produces a *template* .pptx you
     can open in PowerPoint / Keynote / Google Slides and drop the
     remaining PNGs into the placeholder boxes.

  2. Inside Colab after running the main + presentation notebooks —
     finds every figure under
     /content/drive/MyDrive/primates-data/outputs/presentation/ and
     embeds them directly. Produces a *fully populated* .pptx.

The script tries the Drive path first and falls back to the repo path.
Override the Drive path via the PRIMATE_PRESENT_DIR env var.
"""
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# -----------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------
HERE = Path(__file__).parent
REPO_FIGURES = HERE / "figures"
OUT_PATH = HERE / "primate_detection_presentation.pptx"

DRIVE_PRESENT_DIR = Path(os.environ.get(
    "PRIMATE_PRESENT_DIR",
    "/content/drive/MyDrive/primates-data/outputs/presentation",
))


def find_image(name):
    """Find an image by name. Prefer Drive output, fall back to repo figures."""
    for base in (DRIVE_PRESENT_DIR, REPO_FIGURES):
        p = base / name
        if p.exists():
            return p
    return None


# -----------------------------------------------------------------------
# Colors & style
# -----------------------------------------------------------------------
COLOR_TITLE = RGBColor(0x2C, 0x3E, 0x50)     # dark blue-grey
COLOR_BODY = RGBColor(0x34, 0x49, 0x5E)      # slightly lighter
COLOR_MUTED = RGBColor(0x7F, 0x8C, 0x8D)     # grey
COLOR_ACCENT = RGBColor(0x29, 0x80, 0xB9)    # blue
COLOR_PLACEHOLDER = RGBColor(0xBD, 0xC3, 0xC7)


# -----------------------------------------------------------------------
# Speaker notes (kept in sync with presentation_notebooks/speaker_notes.md)
# -----------------------------------------------------------------------
NOTES = {
    1: (
        "(0:10) Hi, I'm [name]. Today I'm going to show you how we used "
        "deep learning to automatically detect primate vocalizations in "
        "multi-hour field recordings from Makokou, Gabon."
    ),
    2: (
        "(0:40) Bioacoustic monitoring — leaving microphones in the forest "
        "— is one of the most powerful tools we have for non-invasive "
        "wildlife monitoring. The problem is scale: a week of continuous "
        "recording produces hundreds of hours of audio, and a human analyst "
        "can listen to maybe one or two hours a day. Our goal is to "
        "automate the task: given a long field recording, find every time "
        "one of three forest primates — Cercopithecus nictitans, "
        "Colobus guereza, and Pan troglodytes — vocalizes, and classify "
        "which species it is."
    ),
    3: (
        "(0:45) Our training set comes from pre-extracted 5-second clips "
        "of each species plus a background class that contains forest "
        "ambient noise and non-target primate species. Around 870 species "
        "clips and 1,800 background clips in total. On the right you can "
        "see one mel-spectrogram per species — this is what the model "
        "actually sees."
    ),
    4: (
        "(1:00) The core idea is simple: treat audio classification as "
        "image classification. We take a 5-second sliding window, convert "
        "it to a mel-spectrogram, resize to 224x224, and feed it to a "
        "VGG19 network pre-trained on ImageNet. The low-level features "
        "VGG19 learned on natural images — edges, textures — are exactly "
        "what distinguishes spectrograms of different calls. We freeze "
        "the backbone and train only a small custom head. Augmentation "
        "(background mixing at varying SNR, time / frequency cropping) "
        "pushes the 870 clips to ~6,000 effective training samples."
    ),
    5: (
        "(1:00) On a held-out validation set the model reaches 94.3% "
        "accuracy across the four classes. The confusion matrix tells a "
        "more nuanced story — look at the diagonal: Cercopithecus and "
        "Colobus recall is above 90%. The hardest class is the background "
        "class, where some samples get mis-labelled as Cercopithecus — "
        "which makes sense because many background clips contain distantly "
        "calling monkeys."
    ),
    6: (
        "(1:10) Validation accuracy on clean clips is not what matters — "
        "what matters is whether it works on real multi-hour field audio. "
        "We ran the model on 13 recordings from Makokou on June 9 2022, "
        "morning to night. Practical surprise: the default confidence "
        "threshold of 0.7 was way too strict. A threshold sweep showed "
        "the confidence distribution is bimodal — the model is either "
        "very sure or it assigns the window to background — so we "
        "settled on 0.4, below which the detection count saturates. "
        "That gave us 44 detections across the 13 files. Parsing the "
        "timestamp out of each filename lets us plot detections by hour "
        "of day, which recovers meaningful diurnal activity patterns. "
        "(Play a detected clip if the room has audio.)"
    ),
    7: (
        "(0:20) Three takeaways: VGG19 transfer learning works well for "
        "primate vocal classification (94% on clean clips); deployed "
        "end-to-end on real 2022 field recordings it recovers meaningful "
        "diurnal activity patterns; and the whole pipeline is released "
        "as a reusable Python package — any researcher with their own "
        "primate recordings can run it by editing one config file. "
        "Main open challenge: Pan troglodytes recall. Next steps: more "
        "field data, per-call-type classification, temporal context. "
        "Thank you — happy to take questions."
    ),
    8: (
        "(backup, only if asked) The whole pipeline lives in src/ as 8 "
        "Python modules — one per pipeline stage. config.py is the entry "
        "point: all paths, hyper-parameters, and the list of species "
        "folders live there. To run this on a new dataset the only file "
        "you edit is config.py — everything else is dataset-agnostic. "
        "Don't read the table row by row; let the audience read it."
    ),
}


# -----------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------
def add_text_box(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=COLOR_BODY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, *,
                size=18, color=COLOR_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + text
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_image_or_placeholder(slide, filename, left, top, width, height, *,
                             caption=None):
    path = find_image(filename)
    if path is not None:
        # Keep aspect ratio, center inside the target box
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        # If it ended up taller than our box, re-add with height constraint
        if pic.height > height:
            slide.shapes._spTree.remove(pic._element)
            pic = slide.shapes.add_picture(str(path), left, top, height=height)
        # Center horizontally inside the box
        pic.left = left + (width - pic.width) // 2
        pic.top = top + (height - pic.height) // 2
        return pic
    # Placeholder rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.background()
    line = rect.line
    line.color.rgb = COLOR_PLACEHOLDER
    line.width = Pt(1.5)
    line.dash_style = 7  # round-dot
    tf = rect.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ drop in: {filename} ]"
    if caption:
        run.text += f"\n{caption}"
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = COLOR_MUTED
    return rect


def set_notes(slide, text):
    nts = slide.notes_slide.notes_text_frame
    nts.text = text


def title(slide, text, *, size=32, color=COLOR_TITLE):
    add_text_box(
        slide, Inches(0.6), Inches(0.35),
        Inches(12.0), Inches(0.9),
        text, size=size, bold=True, color=color,
    )


# -----------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # Blank

    # --- Slide 1: title ---
    s = prs.slides.add_slide(blank)
    add_text_box(
        s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.6),
        "Detecting primate vocalizations\nin long field recordings",
        size=44, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.9),
        "VGG19 transfer learning on mel-spectrograms",
        size=24, color=COLOR_ACCENT, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.6),
        "[Your name]   ·   [Your affiliation]   ·   [Date]",
        size=18, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
    )
    set_notes(s, NOTES[1])

    # --- Slide 2: the problem ---
    s = prs.slides.add_slide(blank)
    title(s, "The problem")
    add_bullets(
        s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.0),
        [
            "Bioacoustic monitoring: leave microphones in the forest, "
            "non-invasive wildlife monitoring, runs 24 / 7 for weeks.",
            "The bottleneck is analysis — a week of continuous recording "
            "can produce hundreds of hours of audio, far more than a "
            "human can listen to.",
            "Goal: automatically find and classify vocalizations of three "
            "forest primates — Cercopithecus nictitans, Colobus guereza, "
            "and Pan troglodytes — in multi-hour Makokou recordings.",
        ],
        size=20,
    )
    set_notes(s, NOTES[2])

    # --- Slide 3: the data ---
    s = prs.slides.add_slide(blank)
    title(s, "Training data: 3 species + background")
    add_image_or_placeholder(
        s, "01_clips_per_class.png",
        Inches(0.5), Inches(1.5), Inches(6.3), Inches(5.6),
    )
    add_image_or_placeholder(
        s, "03_example_spectrograms.png",
        Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.6),
        caption="(one mel-spectrogram per species)",
    )
    add_text_box(
        s, Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.4),
        "~870 species clips (augmented ×7)  +  ~1800 background clips  "
        "=  ~6000 effective training samples",
        size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
    )
    set_notes(s, NOTES[3])

    # --- Slide 4: the method (architecture diagram) ---
    s = prs.slides.add_slide(blank)
    title(s, "Audio-as-image classification via transfer learning")
    add_image_or_placeholder(
        s, "pipeline_architecture.png",
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.2),
    )
    add_bullets(
        s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1.6),
        [
            "5 s sliding window → 128-bin mel-spectrogram → 224×224 RGB",
            "VGG19 ImageNet backbone (frozen) + custom classifier head",
            "Augmentation: background mixing (SNR -5..10 dB), time/freq crop, "
            "freq translation — ×7 effective training-set multiplier",
        ],
        size=15,
    )
    set_notes(s, NOTES[4])

    # --- Slide 5: model results ---
    s = prs.slides.add_slide(blank)
    title(s, "94.3% validation accuracy")
    add_image_or_placeholder(
        s, "02_confusion_matrix.png",
        Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2),
    )
    add_text_box(
        s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.5),
        "Per-class F1 > 0.9 for all three species; hardest class is "
        "background ↔ Cercopithecus (distant calls).",
        size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
    )
    set_notes(s, NOTES[5])

    # --- Slide 6: field deployment ---
    s = prs.slides.add_slide(blank)
    title(s, "Deployed on 13 real recordings from Makokou")
    add_image_or_placeholder(
        s, "05_detections_by_hour.png",
        Inches(0.8), Inches(1.5), Inches(8.2), Inches(5.5),
    )
    add_bullets(
        s, Inches(9.2), Inches(1.7), Inches(4.0), Inches(5.0),
        [
            "13 recordings,\nJune 9 2022,\nmorning → night",
            "Default threshold 0.7 → 18 detections (over-strict)",
            "Threshold sweep → 0.4 = knee of curve",
            "Final: 44 detections,\nbalanced across species",
            "Diurnal activity patterns recovered from filename timestamps",
        ],
        size=13,
    )
    add_text_box(
        s, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.4),
        "(Play one of the extracted clips from outputs/detected_clips/ if the room has audio.)",
        size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
    )
    set_notes(s, NOTES[6])

    # --- Slide 7: takeaways ---
    s = prs.slides.add_slide(blank)
    title(s, "Takeaways")
    add_bullets(
        s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.3),
        [
            "VGG19 transfer learning works for primate vocal classification "
            "— 94.3% validation accuracy.",
            "Deployed end-to-end on real 2022 Makokou recordings, "
            "recovers meaningful diurnal vocal activity patterns.",
            "Released as a reusable Python package — any researcher with "
            "their own primate recordings can run it by editing one config file.",
        ],
        size=22,
    )
    add_text_box(
        s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.7),
        "Open challenge: Pan troglodytes recall. "
        "Next steps: more field data · per-call-type classification · temporal context.",
        size=16, color=COLOR_MUTED,
    )
    add_text_box(
        s, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.6),
        "Thank you — questions welcome.",
        size=20, bold=True, color=COLOR_ACCENT,
    )
    set_notes(s, NOTES[7])

    # --- Slide 8: BACKUP — package structure (hidden) ---
    s = prs.slides.add_slide(blank)
    title(s, "Backup — package structure", color=COLOR_MUTED)
    add_image_or_placeholder(
        s, "package_structure.png",
        Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6),
    )
    add_text_box(
        s, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.4),
        "Hidden from normal flow. Pull up only if a question asks about code / reproducibility.",
        size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Mark slide 8 as hidden
    s_element = s._element
    show = s_element.get("show")
    s_element.set("show", "0")
    set_notes(s, NOTES[8])

    prs.save(OUT_PATH)
    print(f"Saved {OUT_PATH}")

    # Summary of which images were embedded vs placeholdered
    print("\nImage status:")
    wanted = [
        "01_clips_per_class.png",
        "03_example_spectrograms.png",
        "pipeline_architecture.png",
        "02_confusion_matrix.png",
        "05_detections_by_hour.png",
        "package_structure.png",
    ]
    for name in wanted:
        p = find_image(name)
        if p:
            print(f"  OK        {name:40s}  ->  {p}")
        else:
            print(f"  placeholder  {name:40s}  (drop in manually)")


if __name__ == "__main__":
    build()
